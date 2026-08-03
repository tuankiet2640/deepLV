"""Format-preservation tests.

The contract under test is that a translated document is the original document
with different text: styling, structure and everything non-textual come through
untouched. Each format is checked by translating with a marker function and then
re-opening the output with its real library to confirm the formatting is still
attached to the right words.
"""

import io

import pytest

from src.api.services.document_format import (
    SUPPORTED_FORMATS,
    DocumentFormatError,
    file_extension,
    get_handler,
    is_translatable,
)
from src.api.services.document_format.pdf_handler import _split_blocks
from src.api.services.document_format.redistribute import distribute, group_slots
from src.api.services.document_format.segments import split_outer_whitespace
from tests.document_fixtures import (
    HTML_SOURCE,
    TXT_SOURCE,
    build_docx,
    build_pdf,
    build_pptx,
    build_xlsx,
)


def translate_all(handler, data, transform=str.upper):
    """Rebuild a document with every translatable segment transformed."""
    segments = handler.extract(data)
    translations = {s.id: transform(s.text) for s in segments if s.translatable}
    return handler.rebuild(data, translations), segments


def assert_identity_round_trip(handler, data):
    """Rebuilding with the source text must not alter the document's text."""
    segments = handler.extract(data)
    rebuilt = handler.rebuild(data, {s.id: s.text for s in segments})
    again = handler.extract(rebuilt)
    assert [s.text for s in again] == [s.text for s in segments]
    assert [s.context for s in again] == [s.context for s in segments]


# --------------------------------------------------------------------------
# Redistribution: the algorithm that puts translated text back into runs
# --------------------------------------------------------------------------


def test_distribute_is_lossless():
    text = "Der schnelle braune Fuchs springt uber den faulen Hund."
    for weights in ([4, 11, 29], [1, 1], [50, 1, 1, 1], [0, 0, 5]):
        pieces = distribute(text, weights)
        assert len(pieces) == len(weights)
        assert "".join(pieces) == text


def test_distribute_single_slot_keeps_text_whole():
    assert distribute("Hello there", [7]) == ["Hello there"]


def test_distribute_snaps_to_word_boundaries():
    # Equal weights would cut mid-word at index 11; the nearest space wins.
    first, second = distribute("alpha beta gamma delta", [1, 1])
    assert first.strip() and second.strip()
    assert not first.endswith("g")
    assert "".join([first, second]) == "alpha beta gamma delta"


def test_distribute_handles_no_weight_information():
    assert distribute("text", [0, 0]) == ["text", ""]


def test_distribute_without_spaces_still_splits():
    # CJK has no spaces to snap to, so the proportional cut is used as-is.
    pieces = distribute("速い茶色のキツネ", [1, 1])
    assert "".join(pieces) == "速い茶色のキツネ"
    assert all(pieces)


def test_group_slots_merges_only_neighbours():
    assert group_slots(["a", "a", "b", "a"]) == [[0, 1], [2], [3]]
    assert group_slots([]) == []


# --------------------------------------------------------------------------
# Segment filtering
# --------------------------------------------------------------------------


@pytest.mark.parametrize(
    "text",
    ["Hello", "Q1 2024", "速い茶色", "a b", "Total: 5 items"],
)
def test_translatable_text_is_kept(text):
    assert is_translatable(text)


@pytest.mark.parametrize(
    "text",
    [
        "",
        "   ",
        "42",
        "1,240,000",
        "2024-01-15",
        "12.5%",
        "---",
        "https://example.com",
        "sales@example.com",
        "{placeholder}",
    ],
)
def test_untranslatable_text_is_skipped(text):
    assert not is_translatable(text)


def test_split_outer_whitespace_round_trips():
    for value in ["  padded  ", "none", "\n\tlead", "trail \n", "   "]:
        lead, core, trail = split_outer_whitespace(value)
        assert lead + core + trail == value


# --------------------------------------------------------------------------
# DOCX
# --------------------------------------------------------------------------


def test_docx_extracts_every_text_location():
    segments = get_handler("r.docx").extract(build_docx())
    texts = [s.text for s in segments]

    assert "Quarterly Report" in texts
    # Runs are merged into one unit so the provider sees a whole sentence.
    assert "The quick brown fox jumps over the lazy dog." in texts
    # A manual break splits the paragraph rather than swallowing the break.
    assert "Line one" in texts and "Line two" in texts
    # A tab separates two units so it cannot be moved by translation.
    assert "Name:" in texts and "Value here" in texts
    assert "Region" in texts and "Europe" in texts
    assert "Confidential header" in texts
    assert "Page footer text" in texts

    skipped = {s.text for s in segments if not s.translatable}
    assert skipped == {"1,240,000", "42"}


def test_docx_keeps_inline_formatting_on_the_right_words():
    from docx import Document

    handler = get_handler("r.docx", target_lang="de")
    output, _ = translate_all(handler, build_docx())
    doc = Document(io.BytesIO(output))

    body = [p for p in doc.paragraphs if p.text.strip()]
    assert body[0].style.name == "Heading 1"
    assert body[0].text == "QUARTERLY REPORT"

    # Three runs in, three runs out, with bold and colour still on the middle one.
    sentence = body[1]
    assert len(sentence.runs) == 3
    assert sentence.runs[1].bold is True
    assert str(sentence.runs[1].font.color.rgb) == "C00000"
    assert sentence.text == "THE QUICK BROWN FOX JUMPS OVER THE LAZY DOG."


def test_docx_preserves_structure_and_untranslatables():
    from docx import Document

    output, _ = translate_all(get_handler("r.docx", target_lang="de"), build_docx())
    doc = Document(io.BytesIO(output))

    rows = [[cell.text for cell in row.cells] for row in doc.tables[0].rows]
    assert rows == [["REGION", "REVENUE"], ["EUROPE", "1,240,000"]]

    assert doc.sections[0].header.paragraphs[0].text == "CONFIDENTIAL HEADER"
    assert doc.sections[0].footer.paragraphs[0].text == "PAGE FOOTER TEXT"

    # The tab and the manual break survive as layout, not as translated text.
    assert any("\t" in p.text for p in doc.paragraphs)
    assert any("\n" in p.text for p in doc.paragraphs)


def test_docx_identity_round_trip():
    assert_identity_round_trip(get_handler("r.docx", target_lang="de"), build_docx())


def test_docx_keeps_images_and_relationships():
    # The output is the input archive with edited XML, so every part comes along.
    from docx import Document

    source = build_docx()
    output, _ = translate_all(get_handler("r.docx"), source)
    before = {str(p.partname) for p in Document(io.BytesIO(source)).part.package.iter_parts()}
    after = {str(p.partname) for p in Document(io.BytesIO(output)).part.package.iter_parts()}
    assert before == after


# --------------------------------------------------------------------------
# PPTX
# --------------------------------------------------------------------------


def test_pptx_covers_slides_tables_and_notes():
    segments = get_handler("d.pptx").extract(build_pptx())
    contexts = {s.context for s in segments}
    texts = [s.text for s in segments]

    assert "Market Overview" in texts
    assert "Mixed bold tail" in texts
    assert "slide 1 table" in contexts
    assert "slide 1 notes" in contexts
    assert [s.text for s in segments if not s.translatable] == ["31%"]


def test_pptx_keeps_run_formatting_and_notes():
    from pptx import Presentation

    output, _ = translate_all(get_handler("d.pptx", target_lang="de"), build_pptx())
    slide = Presentation(io.BytesIO(output)).slides[0]

    assert slide.shapes.title.text == "MARKET OVERVIEW"

    mixed = slide.placeholders[1].text_frame.paragraphs[1]
    # The fixture sets bold explicitly on every run, so all three keep an
    # explicit value and the middle one is still the bold slot.
    assert [r.font.bold for r in mixed.runs] == [False, True, False]
    assert "".join(r.text for r in mixed.runs) == "MIXED BOLD TAIL"

    assert slide.notes_slide.notes_text_frame.text == "REMEMBER TO MENTION GROWTH"


def test_pptx_identity_round_trip():
    assert_identity_round_trip(get_handler("d.pptx", target_lang="de"), build_pptx())


# --------------------------------------------------------------------------
# XLSX
# --------------------------------------------------------------------------


def test_xlsx_translates_strings_only():
    segments = get_handler("s.xlsx").extract(build_xlsx())
    assert [s.text for s in segments] == ["Product name", "Units sold", "Widget", "Total"]
    assert [s.context for s in segments][0] == "Sales!A1"


def test_xlsx_keeps_styles_numbers_and_formulas():
    from openpyxl import load_workbook

    output, _ = translate_all(get_handler("s.xlsx", target_lang="de"), build_xlsx())
    sheet = load_workbook(io.BytesIO(output))["Sales"]

    assert sheet["A1"].value == "PRODUCT NAME"
    assert sheet["A1"].font.bold is True
    assert sheet["A1"].font.size == 14
    assert sheet["A1"].fill.fgColor.rgb == "00FFFF00"
    assert sheet["B2"].value == 1500
    assert sheet["B3"].value == "=SUM(B2:B2)"
    assert sheet.column_dimensions["A"].width == 32


def test_xlsx_identity_round_trip():
    assert_identity_round_trip(get_handler("s.xlsx"), build_xlsx())


# --------------------------------------------------------------------------
# HTML
# --------------------------------------------------------------------------


def test_html_translates_text_and_visible_attributes():
    segments = get_handler("p.html").extract(HTML_SOURCE)
    by_context = {s.context: s.text for s in segments}

    assert by_context["title"] == "Product page"
    assert by_context["h1"] == "Welcome aboard"
    assert by_context["strong"] == "important"
    assert by_context["a[title]"] == "Go to docs"
    assert by_context["img[alt]"] == "A red bicycle"
    # Script contents are code, never prose.
    assert not any("do not translate" in s.text for s in segments)


def test_html_keeps_markup_and_sets_language():
    output, _ = translate_all(get_handler("p.html", target_lang="de"), HTML_SOURCE)
    text = output.decode("utf-8")

    assert text.startswith("<!DOCTYPE html>")
    assert '<html lang="de">' in text
    assert 'class="hero"' in text
    assert 'href="/docs"' in text
    assert 'src="bike.png"' in text
    assert "<style>.hero{color:red}</style>" in text
    assert 'var keep = "do not translate";' in text
    # Spacing around inline elements is not swallowed by the translation.
    assert "THIS IS <strong>IMPORTANT</strong> TEXT WITH A" in text


def test_html_identity_round_trip():
    assert_identity_round_trip(get_handler("p.html"), HTML_SOURCE)


# --------------------------------------------------------------------------
# TXT
# --------------------------------------------------------------------------


def test_txt_preserves_line_structure_and_endings():
    output, _ = translate_all(get_handler("n.txt"), TXT_SOURCE)
    assert output.decode("utf-8") == (
        "HELLO WORLD\r\n\r\n    INDENTED LINE\r\nSECOND LINE\r\n42\r\n"
    )


def test_txt_falls_back_to_utf8_when_source_encoding_cannot_encode():
    handler = get_handler("n.txt")
    source = "cafe\n".encode("latin-1")
    output, _ = translate_all(handler, source, transform=lambda _: "咖啡")
    assert output.decode("utf-8") == "咖啡\n"


def test_txt_identity_round_trip():
    assert_identity_round_trip(get_handler("n.txt"), TXT_SOURCE)


# --------------------------------------------------------------------------
# PDF
# --------------------------------------------------------------------------


def test_pdf_returns_docx_and_reports_layout_is_not_preserved():
    handler = get_handler("report.pdf", target_lang="fr")
    assert handler.preserves_format is False
    assert handler.output_extension == ".docx"
    assert handler.output_filename("report.pdf") == "report-fr.docx"


def test_pdf_keeps_page_boundaries():
    from docx import Document

    pdf = build_pdf([["Page one text"], ["Page two text"]])
    handler = get_handler("report.pdf", target_lang="fr")
    output, segments = translate_all(handler, pdf)

    assert [s.context for s in segments] == ["page 1", "page 2"]
    paragraphs = Document(io.BytesIO(output)).paragraphs
    assert [p.text for p in paragraphs] == ["PAGE ONE TEXT", "PAGE TWO TEXT"]
    # The second page starts on a new page rather than after a blank paragraph.
    assert paragraphs[1].paragraph_format.page_break_before is True


def test_pdf_without_text_is_rejected():
    handler = get_handler("scan.pdf")
    with pytest.raises(DocumentFormatError, match="OCR"):
        handler.extract(build_pdf([[]]))


def test_pdf_block_splitting_rejoins_wrapped_lines():
    # A hyphenated wrap becomes one word, and a blank line ends the block.
    assert _split_blocks("sever-\nal lines here") == ["several lines here"]
    assert _split_blocks("First para\n\nSecond para") == ["First para", "Second para"]
    # Wrapped lines are rejoined so the provider sees a whole sentence.
    assert _split_blocks("one\ntwo") == ["one two"]
    # A list marker starts its own block.
    assert _split_blocks("Intro line\n- item one\n- item two") == [
        "Intro line",
        "- item one",
        "- item two",
    ]


# --------------------------------------------------------------------------
# Registry
# --------------------------------------------------------------------------


def test_supported_formats_cover_the_office_suite():
    assert SUPPORTED_FORMATS == {"docx", "pptx", "xlsx", "pdf", "html", "htm", "txt"}


def test_unsupported_extension_is_rejected_with_guidance():
    with pytest.raises(DocumentFormatError, match="Unsupported file format"):
        get_handler("archive.zip")


@pytest.mark.parametrize(
    ("filename", "expected"),
    [("a.DOCX", "docx"), ("a.tar.gz", "gz"), ("noextension", ""), ("a.HTM", "htm")],
)
def test_file_extension_normalisation(filename, expected):
    assert file_extension(filename) == expected


def test_office_formats_round_trip_to_their_own_type():
    for filename in ("a.docx", "a.pptx", "a.xlsx", "a.html", "a.txt"):
        handler = get_handler(filename, target_lang="es")
        assert handler.preserves_format is True
        assert handler.output_extension == f".{file_extension(filename)}"
