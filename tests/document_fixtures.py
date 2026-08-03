"""Builders for real documents to exercise the format handlers against.

Everything here produces genuine files in memory, so the tests assert on what the
parsing libraries actually see rather than on hand-written XML that may not match
what Word, PowerPoint or Excel emit.
"""

import io


def build_docx() -> bytes:
    """A Word document with the features that make format preservation hard."""
    from docx import Document
    from docx.shared import RGBColor

    doc = Document()
    doc.add_heading("Quarterly Report", level=1)

    # One sentence split across three runs with different formatting.
    paragraph = doc.add_paragraph()
    paragraph.add_run("The ")
    emphasised = paragraph.add_run("quick brown")
    emphasised.bold = True
    emphasised.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)
    paragraph.add_run(" fox jumps over the lazy dog.")

    # A manual line break must stay put and split the text into two units.
    broken = doc.add_paragraph("Line one")
    broken.add_run().add_break()
    broken.add_run("Line two")

    # A tab is layout, not text.
    doc.add_paragraph("Name:\tValue here")

    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Revenue"
    table.cell(1, 0).text = "Europe"
    table.cell(1, 1).text = "1,240,000"

    doc.sections[0].header.paragraphs[0].text = "Confidential header"
    doc.sections[0].footer.paragraphs[0].text = "Page footer text"

    doc.add_paragraph("Reach us at sales@example.com")
    doc.add_paragraph("42")

    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def build_pptx() -> bytes:
    """A slide with a title, split runs, a table and speaker notes."""
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Market Overview"

    body = slide.placeholders[1].text_frame
    body.text = "First bullet point"
    paragraph = body.add_paragraph()
    for text, bold in (("Mixed ", False), ("bold", True), (" tail", False)):
        run = paragraph.add_run()
        run.text = text
        run.font.bold = bold

    table = slide.shapes.add_table(2, 2, Inches(1), Inches(4), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "Country"
    table.cell(0, 1).text = "Share"
    table.cell(1, 0).text = "Germany"
    table.cell(1, 1).text = "31%"

    slide.notes_slide.notes_text_frame.text = "Remember to mention growth"

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def build_xlsx() -> bytes:
    """A sheet with styling, a formula, numbers and a custom column width."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Sales"
    sheet["A1"] = "Product name"
    sheet["A1"].font = Font(bold=True, size=14)
    sheet["A1"].fill = PatternFill("solid", fgColor="FFFF00")
    sheet["B1"] = "Units sold"
    sheet["A2"] = "Widget"
    sheet["B2"] = 1500
    sheet["A3"] = "Total"
    sheet["B3"] = "=SUM(B2:B2)"
    sheet.column_dimensions["A"].width = 32

    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


HTML_SOURCE = b"""<!DOCTYPE html>
<html lang="en"><head><meta charset="utf-8"><title>Product page</title>
<style>.hero{color:red}</style></head>
<body>
  <h1 class="hero">Welcome aboard</h1>
  <p>This is <strong>important</strong> text with a <a href="/docs" title="Go to docs">link</a>.</p>
  <img src="bike.png" alt="A red bicycle">
  <script>var keep = "do not translate";</script>
  <p>1234</p>
</body></html>"""

TXT_SOURCE = b"Hello world\r\n\r\n    Indented line\r\nSecond line\r\n42\r\n"


def build_pdf(pages: list[list[str]]) -> bytes:
    """A minimal multi-page PDF with Helvetica text lines.

    Written by hand rather than with a PDF generator so the test suite needs no
    extra dependency to cover the PDF path.
    """
    objects: list[bytes] = []
    kids = " ".join(f"{5 + index * 2} 0 R" for index in range(len(pages)))

    objects.append(b"<< /Type /Catalog /Pages 2 0 R >>")
    objects.append(f"<< /Type /Pages /Kids [{kids}] /Count {len(pages)} >>".encode())
    objects.append(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")
    objects.append(b"<< >>")  # filler so the first page object lands on id 5

    for lines in pages:
        content = ["BT /F1 12 Tf 72 720 Td 14 TL"]
        for line in lines:
            escaped = line.replace("\\", r"\\").replace("(", r"\(").replace(")", r"\)")
            content.append(f"({escaped}) Tj T*")
        content.append("ET")
        stream = "\n".join(content).encode("latin-1")

        objects.append(
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 3 0 R >> >> /Contents "
            + str(len(objects) + 2).encode()
            + b" 0 R >>"
        )
        objects.append(
            b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream"
        )

    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for index, payload in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{index} 0 obj\n".encode() + payload + b"\nendobj\n"

    xref_offset = len(out)
    out += f"xref\n0 {len(objects) + 1}\n".encode()
    out += b"0000000000 65535 f \n"
    for offset in offsets:
        out += f"{offset:010d} 00000 n \n".encode()
    out += (
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
        f"startxref\n{xref_offset}\n%%EOF\n"
    ).encode()
    return bytes(out)
