"""PPTX handler.

Walks the slide tree through python-pptx so each unit can be labelled with the
slide it came from, then reuses the shared OOXML row logic on the underlying
DrawingML paragraphs. Group shapes are recursed into, and tables, speaker notes
and chart titles are covered alongside ordinary shape text.
"""

import io
from collections.abc import Iterator
from typing import Any

import structlog

from src.api.services.document_format.base import (
    DocumentFormatError,
    FormatHandler,
    TextUnit,
)
from src.api.services.document_format.oxml import (
    OoxmlTextRow,
    RowSpec,
    a,
    paragraph_rows,
)

log = structlog.get_logger()

PPTX_SPEC = RowSpec(
    paragraph=a("p"),
    run=a("r"),
    text=a("t"),
    run_props=a("rPr"),
    breaks=frozenset({a("br")}),
    # ``a:fld`` is an auto-generated field such as a slide number; its cached
    # text is recomputed by PowerPoint and must be left alone.
    skip=frozenset({a("fld")}),
    # DrawingML keeps formatting in attributes, so the noise to ignore is there
    # rather than in child elements.
    ignored_prop_attrs=frozenset({"lang", "altLang", "dirty", "err", "smtClean", "noProof"}),
)


class PptxHandler(FormatHandler):
    extensions = ("pptx",)
    output_extension = ".pptx"
    output_mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    def load(self, data: bytes) -> Any:
        from pptx import Presentation

        try:
            return Presentation(io.BytesIO(data))
        except Exception as exc:
            raise DocumentFormatError(f"Could not read PPTX file: {exc}") from exc

    def save(self, document: Any) -> bytes:
        buffer = io.BytesIO()
        document.save(buffer)
        return buffer.getvalue()

    def units(self, document: Any) -> Iterator[TextUnit]:
        for number, slide in enumerate(document.slides, start=1):
            label = f"slide {number}"
            yield from self._shape_units(slide.shapes, label)

            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame
                if notes is not None:
                    yield from self._frame_units(notes, f"{label} notes")

    def _shape_units(self, shapes: Any, label: str) -> Iterator[TextUnit]:
        for shape in shapes:
            # Group shapes nest arbitrarily deep; only they expose a child
            # collection, which is a cheaper test than comparing shape types.
            if hasattr(shape, "shapes"):
                yield from self._shape_units(shape.shapes, label)
                continue

            if getattr(shape, "has_text_frame", False):
                yield from self._frame_units(shape.text_frame, label)

            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        yield from self._frame_units(cell.text_frame, f"{label} table")

            if getattr(shape, "has_chart", False):
                yield from self._chart_units(shape, label)

    def _chart_units(self, shape: Any, label: str) -> Iterator[TextUnit]:
        """Chart titles only.

        Category and series labels live in a cached copy of an embedded workbook.
        Rewriting them means editing that workbook and the cached XML in step;
        until that is handled they are left in the source language rather than
        risking a chart that renders inconsistently.
        """
        try:
            chart = shape.chart
            if chart.has_title:
                yield from self._frame_units(chart.chart_title.text_frame, f"{label} chart")
        except (AttributeError, ValueError, KeyError) as exc:
            log.debug("pptx_chart_title_skipped", slide=label, error=str(exc))

    def _frame_units(self, text_frame: Any, label: str) -> Iterator[TextUnit]:
        for paragraph in text_frame.paragraphs:
            for row in paragraph_rows(paragraph._p, PPTX_SPEC):
                yield OoxmlTextRow(row, label, PPTX_SPEC)
